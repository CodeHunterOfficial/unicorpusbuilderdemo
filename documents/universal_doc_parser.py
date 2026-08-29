# documents/universal_doc_parser.py
import os
import sys
import json
import hashlib
import re
import argparse
import tempfile
import shutil
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, List, Tuple, Any

import yaml
from tqdm import tqdm

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from logger_setup import get_file_logger

logger = get_file_logger("doc_parser", "logs/doc_parser.log")

LIBRARIES = {}

def check_library(name: str, import_path: str) -> bool:
    try:
        __import__(import_path)
        LIBRARIES[name] = True
        return True
    except ImportError:
        LIBRARIES[name] = False
        return False

check_library('pdf', 'PyPDF2')
check_library('docx', 'docx')
check_library('ocr', 'pytesseract')
check_library('pptx', 'pptx')
check_library('xlsx', 'openpyxl')
check_library('epub', 'ebooklib')
check_library('rar', 'rarfile')
check_library('7z', 'py7zr')
check_library('bs4', 'bs4')
check_library('pillow', 'PIL')

if LIBRARIES['pdf']:
    from PyPDF2 import PdfReader
if LIBRARIES['docx']:
    from docx import Document
if LIBRARIES['ocr']:
    import pytesseract
if LIBRARIES['pptx']:
    from pptx import Presentation
if LIBRARIES['xlsx']:
    import openpyxl
if LIBRARIES['epub']:
    from ebooklib import epub
    import ebooklib
if LIBRARIES['bs4']:
    from bs4 import BeautifulSoup
if LIBRARIES['rar']:
    import rarfile
if LIBRARIES['7z']:
    import py7zr
if LIBRARIES['pillow']:
    from PIL import Image


def sha256_hex(text: str, trunc: int = 32) -> str:
    h = hashlib.sha256((text or "").encode("utf-8", errors="ignore")).hexdigest()
    return h[:trunc] if trunc else h


def clean_text(text: str, config: dict) -> str:
    if not text:
        return ""
    if config.get('strip_html', True):
        text = re.sub(r'<[^>]+>', '', text)
    if config.get('normalize_whitespace', True):
        text = re.sub(r'\s+', ' ', text)
    return text.strip()


def load_config(config_path: str) -> dict:
    if os.path.exists(config_path):
        with open(config_path, 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f) or {}
            logger.info(f"Configuration loaded from {config_path}")
            return config
    logger.warning(f"Configuration file not found: {config_path}")
    return {}


def build_extension_map(config: dict) -> Dict[str, str]:
    ext_map = {}
    for category, data in config.get('supported_extensions', {}).items():
        for ext in data.get('extensions', []):
            ext_map[ext.lower()] = data.get('type', category)
    return ext_map


def parse_filename_metadata(filename: str, config: dict) -> Dict[str, Optional[str]]:
    name = os.path.splitext(filename)[0]
    separators = config.get('filename_parsing', {}).get('separators', ['_'])
    for sep in separators:
        if sep in name:
            parts = name.split(sep, 1)
            if config.get('filename_parsing', {}).get('author_first', True):
                return {"author": parts[0].strip(), "title": parts[1].strip()}
            else:
                return {"author": parts[1].strip(), "title": parts[0].strip()}
    return {"author": None, "title": name}


def parse_rubric_from_path(file_path: str, root_dir: str, config: dict) -> Optional[str]:
    depth = config.get('rubric_parsing', {}).get('depth', 0)
    ignore_root = config.get('rubric_parsing', {}).get('ignore_root', True)
    try:
        rel_path = os.path.relpath(os.path.dirname(file_path), root_dir)
        if rel_path == '.' and ignore_root:
            return None
        parts = [p for p in rel_path.split(os.sep) if p]
        if len(parts) > depth:
            return parts[depth]
        return None
    except Exception:
        return None


def detect_language(text: str, config: dict) -> str:
    if not config.get('language', {}).get('auto_detect', True):
        return config.get('language', {}).get('default', 'unknown')
    if not text:
        return config.get('language', {}).get('default', 'unknown')
    cyrillic = len(re.findall(r'[а-яёәөүҗңһӣғқҳ]', text, re.I))
    latin = len(re.findall(r'[a-z]', text, re.I))
    arabic = len(re.findall(r'[\u0600-\u06FF]', text))
    if arabic > cyrillic and arabic > latin:
        return "arabic"
    elif cyrillic > latin:
        return "cyrillic"
    elif latin > 0:
        return "latin"
    return config.get('language', {}).get('default', 'unknown')


def extract_text_txt(file_path: str, config: dict) -> Tuple[str, dict]:
    encodings = config.get('text_reading', {}).get('encodings', ['utf-8'])
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read(), {'encoding': enc}
        except (UnicodeDecodeError, UnicodeError):
            continue
    logger.warning(f"Could not decode text file: {file_path}")
    return "", {'error': 'encoding detection failed'}


def extract_text_pdf(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('pdf'):
        return "", {'error': 'PyPDF2 not installed'}
    try:
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return '\n'.join(pages), {'pages': len(reader.pages)}
    except Exception as e:
        logger.error(f"PDF extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_docx(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('docx'):
        return "", {'error': 'python-docx not installed'}
    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells if cell.text)
                if row_text.strip():
                    paragraphs.append(row_text)
        return '\n'.join(paragraphs), {'paragraphs': len(paragraphs)}
    except Exception as e:
        logger.error(f"DOCX extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_pptx(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('pptx'):
        return "", {'error': 'python-pptx not installed'}
    try:
        prs = Presentation(file_path)
        slides = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]
            if slide_text:
                slides.append(f"=== Slide {slide_num} ===\n" + '\n'.join(slide_text))
        return '\n\n'.join(slides), {'slides': len(prs.slides)}
    except Exception as e:
        logger.error(f"PPTX extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_xlsx(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('xlsx'):
        return "", {'error': 'openpyxl not installed'}
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                values = [str(c) for c in row if c is not None]
                if values:
                    rows.append('\t'.join(values))
            if rows:
                sheets.append(f"=== {sheet_name} ===\n" + '\n'.join(rows))
        return '\n\n'.join(sheets), {'sheets': len(wb.sheetnames)}
    except Exception as e:
        logger.error(f"XLSX extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_image(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('ocr') or not LIBRARIES.get('pillow'):
        return "", {'error': 'pytesseract/Pillow not installed'}
    try:
        img = Image.open(file_path)
        if img.width < config.get('ocr', {}).get('min_width', 100) or \
           img.height < config.get('ocr', {}).get('min_height', 100):
            return "", {'error': 'image too small'}
        languages = config.get('ocr', {}).get('languages', 'rus+eng')
        text = pytesseract.image_to_string(img, lang=languages)
        return text, {'size': f"{img.width}x{img.height}"}
    except Exception as e:
        logger.error(f"Image OCR failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_epub(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('epub') or not LIBRARIES.get('bs4'):
        return "", {'error': 'ebooklib/bs4 not installed'}
    try:
        book = epub.read_epub(file_path)
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text()
                if text.strip():
                    chapters.append(text)
        return '\n\n'.join(chapters), {'chapters': len(chapters)}
    except Exception as e:
        logger.error(f"EPUB extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_fb2(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('bs4'):
        return "", {'error': 'BeautifulSoup not installed'}
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        soup = BeautifulSoup(content, 'xml')
        bodies = soup.find_all('body')
        sections = [p.get_text(strip=True) for body in bodies for p in body.find_all('p') if p.get_text(strip=True)]
        return '\n\n'.join(sections), {'sections': len(sections)}
    except Exception as e:
        logger.error(f"FB2 extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_html(file_path: str, config: dict) -> Tuple[str, dict]:
    if not LIBRARIES.get('bs4'):
        return "", {'error': 'BeautifulSoup not installed'}
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        soup = BeautifulSoup(content, 'html.parser')
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        lines = [line.strip() for line in soup.get_text(separator='\n').splitlines() if line.strip()]
        return '\n'.join(lines), {'html_title': soup.title.string if soup.title else None}
    except Exception as e:
        logger.error(f"HTML extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_text_mobi(file_path: str, config: dict) -> Tuple[str, dict]:
    try:
        if LIBRARIES.get('epub'):
            return extract_text_epub(file_path, config)
    except Exception:
        pass
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
        try:
            text = content.decode('utf-8')
            return clean_text(text, {}), {'format': 'mobi_text'}
        except UnicodeDecodeError:
            text = ''.join(chr(b) for b in content if 32 <= b < 127 or b in (10, 13))
            return text, {'format': 'mobi_binary_extract'}
    except Exception as e:
        logger.error(f"MOBI extraction failed for {file_path}: {e}")
        return "", {'error': str(e)}


def extract_files_from_archive(archive_path: str, config: dict, max_depth: int = 3) -> Tuple[List[str], str]:
    """
    Extract files from archive to a temporary directory.
    Returns tuple: (list of extracted file paths, temp_dir path).
    Temp directory is NOT deleted here — caller is responsible.
    """
    extracted_files = []
    ext = os.path.splitext(archive_path)[1].lower()
    temp_dir = tempfile.mkdtemp(prefix="_extracted_")
    
    try:
        if ext == '.zip':
            import zipfile
            with zipfile.ZipFile(archive_path, 'r') as zf:
                zf.extractall(temp_dir)
        elif ext == '.rar' and LIBRARIES.get('rar'):
            with rarfile.RarFile(archive_path, 'r') as rf:
                rf.extractall(temp_dir)
        elif ext == '.7z' and LIBRARIES.get('7z'):
            with py7zr.SevenZipFile(archive_path, 'r') as szf:
                szf.extractall(temp_dir)
        elif ext in ('.tar', '.gz', '.bz2'):
            import tarfile
            with tarfile.open(archive_path, 'r:*') as tf:
                tf.extractall(temp_dir)
        else:
            logger.warning(f"Unsupported archive format: {ext}")
            shutil.rmtree(temp_dir, ignore_errors=True)
            return [], ""
        
        for root, dirs, files in os.walk(temp_dir):
            for f in files:
                extracted_files.append(os.path.join(root, f))
        
        if max_depth > 1:
            nested = []
            for f in extracted_files[:]:
                if os.path.splitext(f)[1].lower() in ('.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'):
                    nested.append(f)
                    extracted_files.remove(f)
            for nf in nested:
                try:
                    nested_files, _ = extract_files_from_archive(nf, config, max_depth - 1)
                    extracted_files.extend(nested_files)
                except Exception as e:
                    logger.error(f"Nested archive extraction failed {nf}: {e}")
        
        return extracted_files, temp_dir
        
    except Exception as e:
        logger.error(f"Archive extraction failed for {archive_path}: {e}")
        shutil.rmtree(temp_dir, ignore_errors=True)
        return [], ""
    
EXTRACTORS = {
    'text': extract_text_txt,
    'pdf': extract_text_pdf,
    'docx': extract_text_docx,
    'pptx': extract_text_pptx,
    'xlsx': extract_text_xlsx,
    'image': extract_text_image,
    'ebook': extract_text_epub,
}

SPECIAL_HANDLERS = {
    '.html': extract_text_html,
    '.htm': extract_text_html,
    '.fb2': extract_text_fb2,
    '.mobi': extract_text_mobi,
    '.azw': extract_text_mobi,
    '.azw3': extract_text_mobi,
}


def parse_documents(root_dir: str, config: dict, output_file: str = None):
    logger.info(f"Starting document parsing in '{root_dir}'")
    ext_map = build_extension_map(config)
    processing = config.get('processing', {})
    output_cfg = config.get('output', {})
    max_size = processing.get('max_file_size_mb', 100)
    min_length = processing.get('min_text_length', 100)
    skip_empty = processing.get('skip_empty', True)

    if not output_file:
        output_dir = output_cfg.get('directory', 'output/documents')
        os.makedirs(output_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = output_cfg.get('filename_template', 'documents_{timestamp}.jsonl').format(timestamp=timestamp)
        output_file = os.path.join(output_dir, filename)
    else:
        output_path = Path(output_file)
        if output_path.suffix:
            os.makedirs(output_path.parent, exist_ok=True)
        else:
            os.makedirs(output_path, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = output_cfg.get('filename_template', 'documents_{timestamp}.jsonl').format(timestamp=timestamp)
            output_file = os.path.join(str(output_path), filename)

    logger.info(f"Output file: {output_file}")
    print(f"[OUTPUT] Saving to: {output_file}")

    all_files = []
    archive_files = []
    for root, dirs, files in os.walk(root_dir):
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        for f in files:
            if f.startswith('.'):
                continue
            ext = os.path.splitext(f)[1].lower()
            if ext in ext_map:
                file_path = os.path.join(root, f)
                if ext_map[ext] == 'archive':
                    archive_files.append(file_path)
                else:
                    all_files.append(file_path)

    temp_dirs_to_clean = []
    if processing.get('recursive_archives', True) and archive_files:
        max_depth = processing.get('max_archive_depth', 3)
        for archive_path in archive_files:
            try:
                extracted, temp_dir = extract_files_from_archive(archive_path, config, max_depth)
                all_files.extend(extracted)
                if temp_dir:
                    temp_dirs_to_clean.append(temp_dir)
                logger.info(f"Extracted {len(extracted)} files from archive: {archive_path}")
            except Exception as e:
                logger.error(f"Failed to extract archive {archive_path}: {e}")

    logger.info(f"Found {len(all_files)} supported files")
    print(f"[FILES] Files found: {len(all_files)}")

    results = []
    stats = {'total': len(all_files), 'processed': 0, 'skipped_size': 0, 'skipped_empty': 0, 'skipped_error': 0, 'by_type': {}}

    for file_path in tqdm(all_files, desc="[PROCESS] Processing", unit="file"):
        try:
            file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
            if file_size_mb > max_size:
                stats['skipped_size'] += 1
                continue
            ext = os.path.splitext(file_path)[1].lower()
            file_type = ext_map.get(ext, 'text')
            filename = os.path.basename(file_path)
            stats['by_type'][file_type] = stats['by_type'].get(file_type, 0) + 1

            if ext in SPECIAL_HANDLERS:
                text, ext_stats = SPECIAL_HANDLERS[ext](file_path, config)
            else:
                extractor = EXTRACTORS.get(file_type, extract_text_txt)
                text, ext_stats = extractor(file_path, config)
            text = clean_text(text, config.get('text_cleaning', {}))

            if skip_empty and (not text or len(text) < min_length):
                stats['skipped_empty'] += 1
                continue

            meta = parse_filename_metadata(filename, config)
            rubric = parse_rubric_from_path(file_path, root_dir, config)
            language = detect_language(text, config)

            item = {
                "url": file_path,
                "title": meta.get("title") or filename,
                "content": text,
                "excerpt": text[:260] + "..." if len(text) > 260 else text,
                "date": datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat(),
                "author": meta.get("author") or config.get('defaults', {}).get('author'),
                "category": rubric or config.get('defaults', {}).get('category'),
                "site": config.get('defaults', {}).get('site', 'local_documents'),
                "hash": sha256_hex(text),
                "language": language or config.get('defaults', {}).get('language'),
                "source_type": config.get('defaults', {}).get('source_type', 'document'),
                "file_name": filename,
                "file_type": file_type,
                "folder": rubric,
                "text_length": len(text),
                "file_size_kb": round(os.path.getsize(file_path) / 1024, 1),
                "extraction_stats": ext_stats,
                "scraped_at": datetime.now().isoformat(),
            }
            results.append(item)
            stats['processed'] += 1
        except Exception as e:
            stats['skipped_error'] += 1
            error_msg = f"{os.path.basename(file_path)}: {e}"
            logger.error(f"Error processing {file_path}: {e}")
            tqdm.write(f"[ERROR] {error_msg}")
            continue

    if results:
        os.makedirs(os.path.dirname(output_file) if os.path.dirname(output_file) else '.', exist_ok=True)
        with open(output_file, 'w', encoding='utf-8') as f:
            for r in results:
                f.write(json.dumps(r, ensure_ascii=False) + "\n")
        if output_cfg.get('save_all', True):
            all_file = output_file.replace('.jsonl', '_all.jsonl')
            with open(all_file, 'w', encoding='utf-8') as f:
                for r in results:
                    f.write(json.dumps(r, ensure_ascii=False) + "\n")
        logger.info(f"Parsing completed. Processed: {stats['processed']}, skipped empty: {stats['skipped_empty']}, errors: {stats['skipped_error']}")
        print(f"\n{'='*60}")
        print(f"[OK] Parsing completed!")
        print(f"  Processed:     {stats['processed']}")
        print(f"  Skipped empty: {stats['skipped_empty']}")
        print(f"  Skipped large: {stats['skipped_size']}")
        print(f"  Errors:        {stats['skipped_error']}")
        print(f"  Total files:   {stats['total']}")
        if stats['by_type']:
            print(f"\n  By type:")
            for ftype, count in sorted(stats['by_type'].items()):
                print(f"    {ftype}: {count}")
        print(f"\n[SAVED] {output_file}")
        print(f"{'='*60}\n")
    else:
        logger.warning("No documents with text found after processing")
        print("[WARN] No documents with text found")
    return {'stats': stats, 'output_file': output_file, 'results_count': len(results)}


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Universal document parser")
    parser.add_argument('folder', help='Folder with documents')
    parser.add_argument('-c', '--config', default='doc_config.yaml', help='Path to config')
    parser.add_argument('-o', '--output', default=None, help='Output JSONL file')
    args = parser.parse_args()
    logger.info(f"CLI started: folder={args.folder}, config={args.config}, output={args.output}")
    config = load_config(args.config)
    parse_documents(args.folder, config, args.output)